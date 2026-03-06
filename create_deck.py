from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Emu(9144000)
prs.slide_height = Emu(5143500)

SNOW_BLUE = RGBColor(0x29, 0xB5, 0xE8)
SNOW_DARK = RGBColor(0x11, 0x27, 0x4F)
DARK_BG = RGBColor(0x0B, 0x1D, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xC9, 0x8D)
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
MEDIUM_GRAY = RGBColor(0x6B, 0x72, 0x80)
BORDER_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
ACN_PURPLE = RGBColor(0xA1, 0x00, 0xFF)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_with_text(slide, shape_type, left, top, width, height, text,
                        fill_color=None, font_size=12, font_color=WHITE,
                        bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(shape_type, Emu(left), Emu(top), Emu(width), Emu(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_multi_text(shape, lines, font_size=10, font_color=WHITE, bold_first=False, line_spacing=1.2):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Calibri"
        if bold_first and i == 0:
            p.font.bold = True
        p.space_before = Pt(2)
        p.space_after = Pt(2)


def add_table(slide, left, top, width, height, rows, cols, data, header_color=SNOW_BLUE):
    table_shape = slide.shapes.add_table(rows, cols, Emu(left), Emu(top), Emu(width), Emu(height))
    table = table_shape.table
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = data[i][j] if i < len(data) and j < len(data[i]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.name = "Calibri"
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LIGHT_GRAY
    return table_shape


# ============================================================
# SLIDE 1: COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 5143500,
                    "", fill_color=DARK_BG)

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 80000,
                    "", fill_color=SNOW_BLUE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 5063500, 9144000, 80000,
                    "", fill_color=ACCENT_GREEN)

add_text_box(slide, 500000, 1200000, 8144000, 800000,
             "Snowflake x Accenture FY27-28", 36, SNOW_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 500000, 1900000, 8144000, 600000,
             "Path to $1B+ | Sustain & Scale", 28, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 500000, 2500000, 8144000, 400000,
             "From Milestone to Market Dominance", 18, ACCENT_CYAN, False, PP_ALIGN.CENTER)
add_text_box(slide, 500000, 3300000, 8144000, 300000,
             "March 2026 | Confidential", 14, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 2500000, 3000000, 4144000, 30000,
                    "", fill_color=ACCENT_PURPLE)

# ============================================================
# SLIDE 2: CONTENTS / AGENDA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 3000000, 5143500,
                    "", fill_color=DARK_BG)

add_text_box(slide, 200000, 300000, 2600000, 500000,
             "PARTNERSHIP\nBUSINESS PLAN", 22, WHITE, True)
add_text_box(slide, 200000, 900000, 2600000, 3800000,
             "This deck outlines our refreshed strategy to sustain and scale beyond $1B.\n\n"
             "Building on FY26's momentum, we introduce bold plays in Agentic AI, "
             "Industry Data Clouds, and Ecosystem Orchestration.\n\n"
             "Questions? Reach out to:\n(SNOW) #team_accenture_global\n(ACN) Global Alliance Leads",
             10, RGBColor(0xCC, 0xCC, 0xCC), False)

sections = [
    ("01", "FY26 Scorecard & Reflection"),
    ("02", "Market & TAM Evolution"),
    ("03", "Bold Growth Plays for FY27-28"),
    ("04", "Agentic AI & Cortex Strategy"),
    ("05", "Industry Data Cloud Expansion"),
    ("06", "Power of Three 2.0"),
    ("07", "Regional Execution Plans"),
    ("08", "COE & Practice Scaling"),
    ("09", "Revenue Waterfall to $1.5B"),
    ("10", "Governance & Mobilization"),
]
y_start = 250000
for num, title in sections:
    add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3200000, y_start, 700000, 350000,
                        num, fill_color=SNOW_BLUE, font_size=16, font_color=WHITE, bold=True)
    add_text_box(slide, 4050000, y_start + 60000, 4800000, 300000,
                 title, 13, RGBColor(0x33, 0x33, 0x33), True)
    y_start += 440000

# ============================================================
# SLIDE 3: DIVIDER - FY26 SCORECARD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 2000000, 9144000, 80000,
                    "", fill_color=SNOW_BLUE)
add_text_box(slide, 500000, 2200000, 8000000, 700000,
             "FY26 Scorecard & Reflection", 36, WHITE, True, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 2900000, 8000000, 400000,
             "Where We've Been — Lessons, Wins & Gaps", 18, ACCENT_CYAN, False, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 3600000, 4000000, 300000,
             "Section 01", 14, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 4: FY26 RESULTS DASHBOARD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "FY26 Results at a Glance", fill_color=DARK_BG, font_size=22, font_color=WHITE, bold=True)

metrics = [
    ("$812M+", "ACN Bookings FY26", "80% toward $1B goal", SNOW_BLUE),
    ("$72M", "SNOW TACV Achieved", "vs $48M prior plan", ACCENT_GREEN),
    ("218+", "Use Cases Deployed", "Across 86 Diamond accounts", ACCENT_PURPLE),
    ("3,000+", "Certifications", "Net-new SnowPro certs", ACCENT_ORANGE),
]

x_pos = 200000
for value, label, sub, color in metrics:
    add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, 800000, 2050000, 1400000,
                        "", fill_color=color)
    add_text_box(slide, x_pos + 100000, 900000, 1850000, 500000,
                 value, 32, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x_pos + 100000, 1350000, 1850000, 300000,
                 label, 12, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x_pos + 100000, 1650000, 1850000, 300000,
                 sub, 9, RGBColor(0xDD, 0xDD, 0xDD), False, PP_ALIGN.CENTER)
    x_pos += 2200000

add_text_box(slide, 200000, 2500000, 8700000, 300000,
             "Key Wins & Momentum Builders", 16, DARK_BG, True)

wins = [
    ["Region", "Key Win", "Impact", "Lesson Forward"],
    ["AMER", "NatWest SAS Migration", "$46M TACV anchor deal", "Replicate migration factory model"],
    ["AMER", "AIG Data Engineering", "Full Cloudera migration", "Package as horizontal offering"],
    ["EMEA", "BMW + Siemens AI/ML", "Manufacturing AI at scale", "Industry-first GTM validated"],
    ["APAC", "BHP Data Platform", "Resource sector expansion", "Mining vertical playbook needed"],
    ["Federal", "CMS IDE Expansion", "Health data modernization", "Power of 3 with AWS works"],
]
add_table(slide, 200000, 2800000, 8700000, 2000000, 6, 4, wins)

# ============================================================
# SLIDE 5: WHAT WORKED / WHAT DIDN'T
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "FY26 Partnership Retrospective", fill_color=DARK_BG, font_size=22, font_color=WHITE, bold=True)

add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 200000, 750000, 4300000, 550000,
                    "", fill_color=ACCENT_GREEN)
add_text_box(slide, 350000, 780000, 4000000, 200000,
             "WHAT WORKED", 16, WHITE, True)
s = add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 200000, 1350000, 4300000, 3500000,
                        "", fill_color=LIGHT_GRAY)
add_multi_text(s, [
    "Migration Factory model drove 60% of new TACV",
    "Power of 3 (AWS) unlocked 3x larger deal sizes",
    "Diamond account focus: 86 activated accounts",
    "AI Refinery narrative resonated with C-suite",
    "SKO & World Tour events drove strong pipeline",
    "3,000+ certifications built credibility",
    "Industry-first GTM in FSI & Manufacturing",
    "Joint innovation workshops at anchor clients",
], font_size=10, font_color=RGBColor(0x33, 0x33, 0x33))

add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4700000, 750000, 4300000, 550000,
                    "", fill_color=RGBColor(0xEF, 0x44, 0x44))
add_text_box(slide, 4850000, 780000, 4000000, 200000,
             "GAPS TO ADDRESS", 16, WHITE, True)
s = add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4700000, 1350000, 4300000, 3500000,
                        "", fill_color=LIGHT_GRAY)
add_multi_text(s, [
    "SNOW AI messaging lags Databricks with ACN CALs",
    "Only 28% of Diamonds activated — 72% whitespace",
    "EMEA & APJ under-invested vs. Americas",
    "Co-sell ground game inconsistent across regions",
    "Deal & service registration friction persists",
    "H&PS and Resources verticals under-penetrated",
    "Agentic AI positioning not yet differentiated",
    "CSP neutrality perception still a challenge",
], font_size=10, font_color=RGBColor(0x33, 0x33, 0x33))

# ============================================================
# SLIDE 6: DIVIDER - MARKET & TAM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 2000000, 9144000, 80000,
                    "", fill_color=ACCENT_GREEN)
add_text_box(slide, 500000, 2200000, 8000000, 700000,
             "Market & TAM Evolution", 36, WHITE, True, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 2900000, 8000000, 400000,
             "The $500B Data & AI Market — Why the Moment is Now", 18, ACCENT_CYAN, False, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 3600000, 4000000, 300000,
             "Section 02", 14, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 7: TAM & MARKET OPPORTUNITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "The Data & AI Services Market Opportunity", fill_color=DARK_BG,
                    font_size=22, font_color=WHITE, bold=True)

add_text_box(slide, 200000, 700000, 8700000, 400000,
             "Snowflake is targeting a $500B+ Total Addressable Market, with 3-Year CAGR of 35%+ — the highest among public cloud data companies.",
             12, DARK_BG, False)

boxes = [
    ("$500B+", "Total Addressable\nMarket by 2028", SNOW_BLUE),
    ("$4.2B+", "Snowflake Product\nRevenue FY26", ACCENT_GREEN),
    ("35%+", "3-Year Revenue\nCAGR", ACCENT_PURPLE),
    ("$1B+", "ACN-SNOW\nPartnership Target", ACCENT_ORANGE),
]
x = 200000
for val, label, col in boxes:
    add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1200000, 2050000, 1100000,
                        "", fill_color=col)
    add_text_box(slide, x + 80000, 1250000, 1900000, 500000,
                 val, 28, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 80000, 1650000, 1900000, 400000,
                 label, 10, WHITE, False, PP_ALIGN.CENTER)
    x += 2200000

shifts = [
    ["Market Shift", "Impact on Partnership", "Our Play"],
    ["Agentic AI mainstream adoption", "Enterprises need trusted agent orchestration", "AI Refinery 2.0 + Cortex Agents"],
    ["Data Gravity moving to clouds", "Legacy migration wave accelerating", "Migration Factory scales globally"],
    ["AI Governance & regulation", "Data governance = board-level priority", "Snowflake Horizon + ACN compliance"],
    ["Industry-specific AI solutions", "Vertical wins > horizontal features", "6 Industry Data Cloud solutions"],
    ["Multi-cloud becomes default", "CSP-agnostic data platforms win", "Power of 3 across AWS/Azure/GCP"],
]
add_table(slide, 200000, 2600000, 8700000, 2200000, 6, 3, shifts)

# ============================================================
# SLIDE 8: DIVIDER - BOLD GROWTH PLAYS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 2000000, 9144000, 80000,
                    "", fill_color=ACCENT_PURPLE)
add_text_box(slide, 500000, 2200000, 8000000, 700000,
             "Bold Growth Plays for FY27-28", 36, WHITE, True, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 2900000, 8000000, 400000,
             "5 Transformative Bets to Sustain $1B+ and Reach $1.5B", 18, ACCENT_CYAN, False, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 3600000, 4000000, 300000,
             "Section 03", 14, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 9: THE 5 BIG BETS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "5 Bold Growth Plays | FY27-28 Strategy", fill_color=DARK_BG,
                    font_size=22, font_color=WHITE, bold=True)

bets = [
    ("01", "AGENTIC AI\nFACTORY", "Scale AI Refinery 2.0\nwith Cortex Agents\nacross 50+ accounts", SNOW_BLUE),
    ("02", "INDUSTRY DATA\nCLOUD BLITZ", "6 vertical solutions\nwith native apps\n& marketplace", ACCENT_GREEN),
    ("03", "MIGRATION\nMEGA-WAVE", "Target $300M in\nlegacy-to-Snowflake\nmigrations", ACCENT_PURPLE),
    ("04", "DIAMOND\nSATURATION", "Activate 200+ of 307\nDiamond accounts\n(65% penetration)", ACCENT_ORANGE),
    ("05", "ECOSYSTEM\nFLYWHEEL", "Power of 3 v2.0\nwith AWS + Azure\n+ SAP + SFDC", RGBColor(0xEC, 0x48, 0x99)),
]

x = 120000
for num, title, desc, color in bets:
    add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 750000, 1720000, 4100000,
                        "", fill_color=color)
    add_text_box(slide, x + 80000, 800000, 1560000, 300000,
                 num, 28, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 80000, 1200000, 1560000, 600000,
                 title, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 80000, 2000000, 1560000, 2500000,
                 desc, 10, RGBColor(0xEE, 0xEE, 0xEE), False, PP_ALIGN.CENTER)
    x += 1800000

# ============================================================
# SLIDE 10: CREATIVE IDEA 1 - AGENTIC AI FACTORY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Play #1: Agentic AI Factory — Cortex-Powered Enterprise Agents",
                    fill_color=SNOW_BLUE, font_size=18, font_color=WHITE, bold=True)

add_text_box(slide, 200000, 700000, 4200000, 300000,
             "The Vision", 16, DARK_BG, True)
add_text_box(slide, 200000, 950000, 4200000, 800000,
             "Create an industrialized 'Agentic AI Factory' combining Accenture's AI Refinery "
             "with Snowflake Cortex to deliver pre-built, governed AI agents at scale. "
             "Position ACN+SNOW as the only partnership that can deploy trusted AI agents "
             "from data to production in 6 weeks.",
             10, RGBColor(0x33, 0x33, 0x33), False)

add_text_box(slide, 4700000, 700000, 4200000, 300000,
             "Revenue Impact", 16, DARK_BG, True)

impacts = [
    ["Metric", "FY27 Target", "FY28 Target"],
    ["Agent deployments", "50 accounts", "150 accounts"],
    ["Average deal size", "$2.5M", "$4M"],
    ["Total pipeline", "$125M", "$600M"],
    ["TACV contribution", "$15M", "$40M"],
    ["Consumption uplift", "3x per account", "5x per account"],
]
add_table(slide, 4700000, 1000000, 4200000, 1200000, 6, 3, impacts)

add_text_box(slide, 200000, 2000000, 8700000, 300000,
             "Creative Execution Ideas", 16, DARK_BG, True)

ideas = [
    ["Initiative", "Description", "Timeline", "Owner"],
    ["Agent Marketplace", "Pre-built Cortex agents for top 10 use cases (Claims, Supply Chain, Marketing, HR, Finance)", "Q1-Q2 FY27", "COE + Product"],
    ["Agentic AI Hackathons", "Quarterly hack events at Diamond clients: build agents in 48hrs using Cortex", "Ongoing", "Field + Marketing"],
    ["Trusted Agent Certification", "New SnowPro specialty cert for Agentic AI — target 500 ACN consultants", "Q2 FY27", "Enablement"],
    ["C-Suite Agent Experience", "Executive briefing center demo: show a live agent solving their industry problem", "Q1 FY27", "GTM + Solutions"],
    ["Agent ROI Calculator", "Self-service tool: quantify agent value by industry vertical", "Q1 FY27", "Solutions"],
]
add_table(slide, 200000, 2300000, 8700000, 2500000, 6, 4, ideas)

# ============================================================
# SLIDE 11: CREATIVE IDEA 2 - INDUSTRY DATA CLOUD BLITZ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Play #2: Industry Data Cloud Blitz — Vertical-First GTM",
                    fill_color=ACCENT_GREEN, font_size=18, font_color=WHITE, bold=True)

add_text_box(slide, 200000, 700000, 8700000, 400000,
             "Launch 6 Industry Data Cloud solutions as native Snowflake apps, each with a dedicated GTM motion, "
             "anchor customer, and consumption-based pricing. Shift from horizontal to vertical dominance.",
             11, RGBColor(0x33, 0x33, 0x33), False)

solutions = [
    ["Industry Cloud", "Native App", "Anchor Customers", "TAM", "FY27 Target", "FY28 Target"],
    ["FSI Claims Intelligence", "Claims AI 2.0 + Fraud Detection Agent", "AIG, Hartford, Chubb", "$8B", "$25M", "$60M"],
    ["Life Sciences Clinical AI", "CSIG + Trial Analytics + Regulatory AI", "Pfizer, Roche, AstraZeneca", "$6B", "$20M", "$50M"],
    ["Manufacturing Digital Twin", "Shop Floor 2.0 + Predictive Quality", "BMW, Siemens, Caterpillar", "$7B", "$22M", "$55M"],
    ["Retail Relevance Engine", "Personalization + Supply Chain AI", "Mars, Nestle, Marriott", "$5B", "$18M", "$45M"],
    ["Energy Grid Intelligence", "Asset Optimization + Carbon Analytics", "Vistra, EDF, PPL", "$4B", "$12M", "$30M"],
    ["Public Sector GovCloud", "FedGenius + Compliance Automation", "CMS, DoD, US Courts", "$3B", "$10M", "$25M"],
]
add_table(slide, 200000, 1200000, 8700000, 2000000, 7, 6, solutions)

add_text_box(slide, 200000, 3400000, 8700000, 300000,
             "Creative Differentiator: \"Industry in a Box\"", 14, DARK_BG, True)
add_text_box(slide, 200000, 3700000, 8700000, 1200000,
             "Each Industry Data Cloud comes pre-packaged with: Reference Architecture | Demo Environment | "
             "First Call Deck | Battle Card | ROI Model | 3 Pre-built Use Cases | Anchor Customer Story | "
             "Dedicated Industry Squad (ACN + SNOW). Deploy-ready in 4 weeks vs. 12+ weeks today.",
             10, RGBColor(0x33, 0x33, 0x33), False)

# ============================================================
# SLIDE 12: CREATIVE IDEA 3 - MIGRATION MEGA-WAVE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Play #3: Migration Mega-Wave — $300M Legacy Modernization",
                    fill_color=ACCENT_PURPLE, font_size=18, font_color=WHITE, bold=True)

add_text_box(slide, 200000, 700000, 4200000, 300000,
             "The Opportunity", 16, DARK_BG, True)
add_text_box(slide, 200000, 950000, 4200000, 900000,
             "70% of enterprises still run critical workloads on legacy platforms (Teradata, SAS, Hadoop, Oracle, "
             "SAP BW). Our NatWest success proves the Migration Factory model works. "
             "Now scale it globally with AI-powered migration tooling — SnowConvert + Agentic code translation.",
             10, RGBColor(0x33, 0x33, 0x33), False)

add_text_box(slide, 4700000, 700000, 4200000, 300000,
             "Migration Targets by Platform", 16, DARK_BG, True)

migrations = [
    ["Source Platform", "Est. Market", "Target Migrations", "Avg Deal Size"],
    ["Teradata", "$50B installed base", "30 accounts", "$5-15M"],
    ["SAS / SAS Viya", "$3B annual spend", "25 accounts", "$3-8M"],
    ["Hadoop / Cloudera", "$2B market", "20 accounts", "$2-6M"],
    ["Oracle Exadata", "$15B installed", "15 accounts", "$4-10M"],
    ["SAP BW/HANA Analytics", "$8B SAP analytics", "15 accounts", "$3-8M"],
    ["On-prem SQL/Netezza", "$5B legacy DW", "20 accounts", "$2-5M"],
]
add_table(slide, 4700000, 1000000, 4200000, 1600000, 7, 4, migrations)

add_text_box(slide, 200000, 2100000, 8700000, 300000,
             "Creative: \"Zero-Risk Migration Guarantee\"", 14, DARK_BG, True)
add_text_box(slide, 200000, 2400000, 8700000, 400000,
             "Offer a money-back guarantee on Phase 1 migrations (assessment + POC) — funded by Snowflake MDF + ACN investment. "
             "If the client doesn't see 30%+ TCO reduction in the first 90 days, the assessment is free.",
             10, RGBColor(0x33, 0x33, 0x33), False)

creative_migrations = [
    ["Creative Initiative", "Description", "Impact"],
    ["AI-Powered Code Translation", "Use Cortex LLMs to auto-convert SAS/Teradata/Oracle to Snowpark", "80% faster migrations"],
    ["Migration-as-a-Service (MaaS)", "Subscription model: ongoing migration + managed modernization", "Recurring $2-5M/yr per client"],
    ["Executive Migration Summits", "Quarterly C-suite events: \"Legacy is the New Risk\" narrative", "20+ qualified leads per event"],
    ["Migration Leaderboard", "Gamified internal competition across ACN regions", "3x pipeline velocity"],
    ["Partner Migration Fund", "Joint $10M fund for migration POCs and accelerators", "50+ POCs per year"],
]
add_table(slide, 200000, 2900000, 8700000, 1900000, 6, 3, creative_migrations)

# ============================================================
# SLIDE 13: CREATIVE IDEA 4 - DIAMOND SATURATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Play #4: Diamond Saturation — 200+ Accounts Activated",
                    fill_color=ACCENT_ORANGE, font_size=18, font_color=WHITE, bold=True)

add_text_box(slide, 200000, 700000, 8700000, 400000,
             "Only 28% of ACN's 307 Diamond accounts are active with Snowflake today. "
             "The biggest growth lever is activating the remaining 72% — an incremental $400M+ opportunity.",
             11, DARK_BG, False)

diamond_data = [
    ["Region", "Total Diamonds", "Active FY26", "FY27 Target", "FY28 Target", "Gap to Close"],
    ["Americas", "142", "45 (32%)", "85 (60%)", "110 (77%)", "+65"],
    ["EMEA", "108", "22 (20%)", "55 (51%)", "80 (74%)", "+58"],
    ["Asia Pacific", "57", "21 (37%)", "35 (61%)", "45 (79%)", "+24"],
    ["TOTAL", "307", "88 (28%)", "175 (57%)", "235 (77%)", "+147"],
]
add_table(slide, 200000, 1200000, 8700000, 1200000, 6, 6, diamond_data)

add_text_box(slide, 200000, 2600000, 8700000, 300000,
             "Creative Account Activation Playbook", 14, DARK_BG, True)

activation = [
    ["Tactic", "Description", "Expected Result"],
    ["\"First 90 Days\" Blitz Kit", "Pre-packaged activation kit for new accounts: discovery template, use case library, POC framework", "30% faster first engagement"],
    ["Account Twinning Program", "Pair inactive Diamonds with successful reference accounts in same industry", "2x activation rate"],
    ["CAL Incentive Program", "Bonus structure for ACN Client Account Leads who activate Snowflake in their portfolio", "50+ new accounts"],
    ["Data & AI Health Check", "Free 2-week assessment of client's data maturity with Snowflake benchmarks", "75% conversion to paid engagement"],
    ["Executive Sponsor Mapping", "Map every Diamond CEO/CTO/CDAO to ACN+SNOW executive for quarterly touchpoints", "Board-level sponsorship"],
]
add_table(slide, 200000, 2900000, 8700000, 1900000, 6, 3, activation)

# ============================================================
# SLIDE 14: CREATIVE IDEA 5 - ECOSYSTEM FLYWHEEL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Play #5: Ecosystem Flywheel — Power of 3 v2.0",
                    fill_color=RGBColor(0xEC, 0x48, 0x99), font_size=18, font_color=WHITE, bold=True)

add_text_box(slide, 200000, 700000, 8700000, 400000,
             "Evolve from ad-hoc Power of 3 deals to a systematic Ecosystem Flywheel: "
             "every ACN Business Group (AWS, Azure, SAP, SFDC) has embedded Snowflake plays with pre-assembled commercial models.",
             11, RGBColor(0x33, 0x33, 0x33), False)

ecosystem = [
    ["Business Group", "FY26 Status", "FY27-28 Play", "Target Pipeline", "Key Accounts"],
    ["AWS (AABG)", "Aligned on NatWest, DTCC", "Agentic AI + Migration on AWS", "$150M", "NatWest, JPMC, CBA"],
    ["Azure (AMBG)", "Establishing connections", "RCG + Manufacturing on Azure", "$120M", "Caterpillar, Nestle, Ross"],
    ["SAP (SBG)", "Early stage", "SAP BW Migration + Data Mesh", "$100M", "BMW, Siemens, Halliburton"],
    ["Salesforce (SFBG)", "Not started", "C360 + Relevance Engine", "$50M", "Disney, Comcast, AT&T"],
    ["Google Cloud (NEW)", "Greenfield", "Multi-cloud Analytics + AI", "$30M", "Snap, Pinterest, OpenAI"],
]
add_table(slide, 200000, 1200000, 8700000, 1400000, 6, 5, ecosystem)

add_text_box(slide, 200000, 2800000, 8700000, 300000,
             "Creative: \"Ecosystem Multiplier\" Model", 14, DARK_BG, True)

multiplier = [
    ["Initiative", "How It Works", "Revenue Multiplier"],
    ["Pre-Built Commercial Models", "3-way pricing templates pre-approved by legal for each BG", "2x deal velocity"],
    ["Joint Innovation Labs", "Permanent labs in NYC, London, Sydney with ACN+SNOW+CSP", "5x qualified pipeline"],
    ["Marketplace Revenue Share", "Snowflake Marketplace apps with ACN implementation services bundled", "$20M+ new revenue stream"],
    ["BG Seller Incentives", "Commission kicker for BG sellers who include Snowflake in deals", "3x Snowflake inclusion rate"],
    ["Quarterly Ecosystem War Rooms", "Joint pipeline review with all BG leads + Snowflake", "40% pipeline conversion"],
]
add_table(slide, 200000, 3100000, 8700000, 1700000, 6, 3, multiplier)

# ============================================================
# SLIDE 15: DIVIDER - AGENTIC AI & CORTEX
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 2000000, 9144000, 80000,
                    "", fill_color=SNOW_BLUE)
add_text_box(slide, 500000, 2200000, 8000000, 700000,
             "Agentic AI & Cortex Strategy", 36, WHITE, True, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 2900000, 8000000, 400000,
             "Building the Enterprise AI Operating System Together", 18, ACCENT_CYAN, False, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 3600000, 4000000, 300000,
             "Section 04", 14, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 16: CORTEX + AI REFINERY DEEP DIVE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Cortex + AI Refinery 2.0: The Enterprise Agent Platform",
                    fill_color=SNOW_BLUE, font_size=18, font_color=WHITE, bold=True)

add_text_box(slide, 200000, 700000, 4200000, 300000,
             "Platform Architecture", 16, DARK_BG, True)

layers = [
    ("Agent Orchestration Layer", "Multi-agent workflows, task routing, human-in-the-loop", SNOW_BLUE),
    ("Cortex AI Services", "LLM functions, embeddings, search, classification, summarization", ACCENT_GREEN),
    ("Data Foundation", "Governed data, Iceberg, Dynamic Tables, Streams & Tasks", ACCENT_PURPLE),
    ("Security & Governance", "Horizon policies, RBAC, data masking, audit trails", DARK_BG),
]

y = 1050000
for title, desc, color in layers:
    add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 200000, y, 4200000, 700000,
                        "", fill_color=color)
    add_text_box(slide, 300000, y + 50000, 4000000, 250000,
                 title, 12, WHITE, True)
    add_text_box(slide, 300000, y + 320000, 4000000, 350000,
                 desc, 9, RGBColor(0xDD, 0xDD, 0xDD), False)
    y += 780000

add_text_box(slide, 4700000, 700000, 4200000, 300000,
             "Top 10 Agent Use Cases for FY27", 14, DARK_BG, True)

agents = [
    ["#", "Agent Use Case", "Industry", "Est. Deal Size"],
    ["1", "Claims Processing Agent", "Insurance", "$3-5M"],
    ["2", "Clinical Trial Analytics Agent", "Life Sciences", "$2-4M"],
    ["3", "Supply Chain Optimizer Agent", "Manufacturing", "$2-5M"],
    ["4", "Customer 360 Agent", "Retail/CPG", "$1-3M"],
    ["5", "Regulatory Compliance Agent", "FSI/Pharma", "$2-4M"],
    ["6", "Energy Grid Forecasting Agent", "Utilities", "$1-3M"],
    ["7", "HR & Talent Intelligence Agent", "Cross-Industry", "$1-2M"],
    ["8", "Marketing Mix Optimizer Agent", "Retail/Media", "$1-3M"],
    ["9", "Fraud Detection Agent", "Banking", "$3-5M"],
    ["10", "Document Intelligence Agent", "Government", "$1-3M"],
]
add_table(slide, 4700000, 1000000, 4200000, 3800000, 11, 4, agents)

# ============================================================
# SLIDE 17: DIVIDER - REGIONAL EXECUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 2000000, 9144000, 80000,
                    "", fill_color=ACCENT_ORANGE)
add_text_box(slide, 500000, 2200000, 8000000, 700000,
             "Regional Execution Plans", 36, WHITE, True, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 2900000, 8000000, 400000,
             "Americas | EMEA | APJ — Tailored for Local Impact", 18, ACCENT_CYAN, False, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 3600000, 4000000, 300000,
             "Section 07", 14, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 18: REGIONAL TARGETS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Regional Revenue & Growth Targets — FY27-28",
                    fill_color=DARK_BG, font_size=20, font_color=WHITE, bold=True)

regional = [
    ["Region", "FY26 Actual", "FY27 Target", "FY28 Target", "Growth %", "Key Plays"],
    ["Americas (inc. Federal)", "$550M", "$700M", "$900M", "27% / 29%", "Diamond saturation, Migration wave, FSI expansion"],
    ["EMEA", "$180M", "$280M", "$400M", "56% / 43%", "Industry-first GTM, Migration Factory, NatWest replication"],
    ["Asia Pacific", "$82M", "$140M", "$200M", "71% / 43%", "ANZ acceleration, Japan expansion, Mining vertical"],
    ["TOTAL", "$812M", "$1,120M", "$1,500M", "38% / 34%", ""],
]
add_table(slide, 200000, 750000, 8700000, 1200000, 5, 6, regional)

add_text_box(slide, 200000, 2100000, 8700000, 300000,
             "SNOW Metrics Aligned", 14, DARK_BG, True)

snow_targets = [
    ["Metric", "FY26 Actual", "FY27 Target", "FY28 Target"],
    ["TACV", "$72M", "$110M", "$160M"],
    ["Use Case Go-Lives", "218", "350", "500"],
    ["Consumption Revenue", "~$24M", "$45M", "$75M"],
    ["Active Diamond Accounts", "88", "175", "235"],
    ["Net New Certifications", "3,000", "4,000", "5,000"],
    ["Active Certified Practitioners", "5,000", "8,000", "12,000"],
]
add_table(slide, 200000, 2400000, 8700000, 1800000, 7, 4, snow_targets)

add_text_box(slide, 200000, 4400000, 8700000, 500000,
             "Creative: Regional \"Lightning Strike\" Campaigns — Each region gets a quarterly blitz week "
             "with exec sponsorship, dedicated pipeline generation, and account activation sprints.",
             10, RGBColor(0x66, 0x66, 0x66), False)

# ============================================================
# SLIDE 19: AMERICAS DETAILED PLAN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Americas Plan — From $550M to $900M",
                    fill_color=SNOW_BLUE, font_size=20, font_color=WHITE, bold=True)

amer = [
    ["Industry", "FY26 Bookings", "FY27 Target", "FY28 Target", "Top Accounts", "Key Solution"],
    ["Financial Services", "$150M", "$200M", "$270M", "JPMC, AIG, Truist, BNY", "Claims AI, Fraud Agent"],
    ["Comms Media Tech", "$85M", "$120M", "$160M", "Disney, AT&T, Comcast", "Relevance Engine 2.0"],
    ["Products (RCG+LS)", "$120M", "$160M", "$210M", "Pfizer, Marriott, Mars, CAT", "Industry Data Clouds"],
    ["Health & Public Svc", "$60M", "$90M", "$120M", "UHG, Centene, BCBS", "CSIG, Clinical AI"],
    ["Resources", "$55M", "$80M", "$100M", "Vistra, Halliburton, PPL", "Energy Grid Intelligence"],
    ["Federal / SLED", "$30M", "$50M", "$40M", "CMS, DoD, Courts", "GovCloud, FedGenius"],
    ["TOTAL Americas", "$550M", "$700M", "$900M", "", ""],
]
add_table(slide, 100000, 700000, 8944000, 2200000, 8, 6, amer)

add_text_box(slide, 200000, 3100000, 8700000, 300000,
             "Americas Creative Plays", 14, DARK_BG, True)

amer_creative = [
    ["Initiative", "Description", "Q1-Q2", "Q3-Q4"],
    ["Diamond Sprint Program", "Dedicated 8-week activation for top 30 inactive Diamonds", "Wave 1: 15 accounts", "Wave 2: 15 accounts"],
    ["Industry Roadshows", "6-city tour: FSI, MFG, LS, RCG, Energy, Federal", "NYC, Chicago, SF", "Dallas, Atlanta, DC"],
    ["Client Innovation Labs", "Permanent AI lab experiences at 3 locations", "NYC (launch)", "SF, Chicago"],
    ["Executive Dinner Series", "Quarterly C-suite dinners themed around Agentic AI", "4 events", "4 events"],
]
add_table(slide, 200000, 3400000, 8700000, 1500000, 5, 4, amer_creative)

# ============================================================
# SLIDE 20: EMEA DETAILED PLAN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "EMEA Plan — From $180M to $400M (2.2x Growth)",
                    fill_color=ACCENT_GREEN, font_size=20, font_color=WHITE, bold=True)

emea = [
    ["Market Unit", "FY26", "FY27 Target", "FY28 Target", "Must-Win Accounts", "Priority Play"],
    ["UKI", "$55M", "$90M", "$130M", "NatWest, HSBC, AstraZeneca, GSK", "Migration Factory, FSI"],
    ["ASG (DACH)", "$45M", "$75M", "$110M", "Siemens, BMW, Merck, Allianz", "MFG Digital Twin, SAP"],
    ["Gallia (France)", "$40M", "$60M", "$85M", "Credit Agricole, Sanofi, Pernod", "FSI, LS Clinical AI"],
    ["Iberia", "$15M", "$25M", "$35M", "Santander, Mapfre, Amadeus", "Migration, FSI"],
    ["Nordics + Italy", "$25M", "$30M", "$40M", "Ericsson, Stellantis, Nordea", "Telco, Industrial"],
    ["TOTAL EMEA", "$180M", "$280M", "$400M", "", ""],
]
add_table(slide, 100000, 700000, 8944000, 1800000, 7, 6, emea)

add_text_box(slide, 200000, 2700000, 8700000, 300000,
             "EMEA Creative Acceleration Plays", 14, DARK_BG, True)

emea_creative = [
    ["Initiative", "Description", "Impact"],
    ["\"NatWest Playbook\" Replication", "Package NatWest SAS migration as repeatable offering for 10 EMEA FSI accounts", "$50M+ pipeline"],
    ["EMEA AI Governance Summit", "Position ACN+SNOW as EU AI Act compliance leaders with Horizon governance", "C-suite credibility"],
    ["Local Language Agent Demos", "Cortex agents in German, French, Spanish for local market resonance", "3x engagement"],
    ["Pan-European Migration Factory", "Centralized migration COE in Bangalore + local delivery pods", "40% cost reduction"],
    ["SnowDays Europe Tour", "10-city enablement tour targeting ACN delivery leads", "500+ enabled consultants"],
]
add_table(slide, 200000, 3000000, 8700000, 1800000, 6, 3, emea_creative)

# ============================================================
# SLIDE 21: APJ DETAILED PLAN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "APJ Plan — From $82M to $200M (2.4x Growth)",
                    fill_color=ACCENT_PURPLE, font_size=20, font_color=WHITE, bold=True)

apj = [
    ["Sub-Region", "FY26", "FY27 Target", "FY28 Target", "Must-Win Accounts", "Priority"],
    ["ANZ", "$30M", "$55M", "$80M", "BHP, CBA, NAB, Coles, Qantas", "Mining, FSI, Retail"],
    ["Japan", "$25M", "$45M", "$65M", "NTT, Toyota, Sony, Mitsubishi", "MFG, Telco"],
    ["SEA + India", "$18M", "$28M", "$38M", "DBS, Singtel, Tata, Reliance", "FSI, Telco, IT"],
    ["Greater China", "$9M", "$12M", "$17M", "Lenovo, Huawei", "Tech, MFG"],
    ["TOTAL APJ", "$82M", "$140M", "$200M", "", ""],
]
add_table(slide, 100000, 700000, 8944000, 1400000, 6, 6, apj)

add_text_box(slide, 200000, 2300000, 8700000, 300000,
             "APJ Creative Growth Plays", 14, DARK_BG, True)

apj_creative = [
    ["Initiative", "Description", "Impact"],
    ["Mining Data Cloud", "Purpose-built solution for BHP/Rio Tinto: ore grade prediction, safety analytics", "New $50M vertical"],
    ["Japan Enterprise Blitz", "Dedicated Japanese-speaking team, localized solutions, keiretsu-model selling", "2x Japan revenue"],
    ["APJ Gen AI Studio", "Permanent innovation centers in Sydney, Tokyo, Singapore", "Hands-on client experiences"],
    ["Cross-border Data Sharing", "Snowflake data sharing for APAC supply chains (ASEAN corridor)", "Unique value prop"],
    ["Accenture Song + Snowflake", "Combined marketing analytics for APAC retail brands", "$15M new pipeline"],
]
add_table(slide, 200000, 2600000, 8700000, 1800000, 6, 3, apj_creative)

# ============================================================
# SLIDE 22: DIVIDER - COE & PRACTICE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 2000000, 9144000, 80000,
                    "", fill_color=ACCENT_GREEN)
add_text_box(slide, 500000, 2200000, 8000000, 700000,
             "COE & Practice Scaling", 36, WHITE, True, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 2900000, 8000000, 400000,
             "Building the World's Largest Snowflake Practice — 15,000+ Practitioners", 18, ACCENT_CYAN, False, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 3600000, 4000000, 300000,
             "Section 08", 14, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 23: COE EVOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Global COE Evolution — From Practice to Business Group",
                    fill_color=DARK_BG, font_size=20, font_color=WHITE, bold=True)

add_text_box(slide, 200000, 700000, 8700000, 400000,
             "The Snowflake Practice is on track to become an Accenture Business Group in CY2027 — "
             "joining AWS, Azure, SAP, and Salesforce as a top-5 technology alliance.",
             11, DARK_BG, False)

coe = [
    ["Pillar", "FY26 State", "FY27 Target", "FY28 Target", "Investment Required"],
    ["Certified Practitioners", "5,000 active", "8,000 active", "12,000 active", "Training programs + cert subsidies"],
    ["Delivery Centers", "3 (US, India, UK)", "6 (+Germany, Australia, Japan)", "8 (+Singapore, Brazil)", "$5M infra investment"],
    ["Industry Solutions", "7 solutions", "12 solutions", "18 solutions", "$8M solution development"],
    ["Full-time COE Staff", "15 (global)", "35 (global)", "60 (global)", "$12M headcount"],
    ["Revenue per Head", "$160K", "$200K", "$250K", "Efficiency + automation"],
    ["Agentic AI Specialists", "0 dedicated", "50 certified", "200 certified", "New cert track needed"],
]
add_table(slide, 200000, 1200000, 8700000, 1800000, 7, 5, coe)

add_text_box(slide, 200000, 3200000, 8700000, 300000,
             "Creative: \"Snowflake Academy\" Program", 14, DARK_BG, True)
add_text_box(slide, 200000, 3500000, 8700000, 1400000,
             "Launch a dedicated Snowflake Academy within Accenture:\n"
             "  - 12-week immersive bootcamp for top talent (200 graduates per quarter)\n"
             "  - Specialization tracks: Data Engineering, AI/ML, Migrations, Industry Solutions\n"
             "  - \"Earn while you learn\" model with guaranteed project placement\n"
             "  - Joint faculty: Snowflake engineers + Accenture delivery leaders\n"
             "  - Certification guarantee: 95%+ pass rate target\n"
             "  - Alumni network becomes the backbone of the global practice",
             10, RGBColor(0x33, 0x33, 0x33), False)

# ============================================================
# SLIDE 24: DIVIDER - REVENUE WATERFALL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 2000000, 9144000, 80000,
                    "", fill_color=ACCENT_ORANGE)
add_text_box(slide, 500000, 2200000, 8000000, 700000,
             "Revenue Waterfall to $1.5B", 36, WHITE, True, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 2900000, 8000000, 400000,
             "The Math Behind Sustaining and Doubling the $1B Milestone", 18, ACCENT_CYAN, False, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 3600000, 4000000, 300000,
             "Section 09", 14, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 25: REVENUE BRIDGE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Revenue Bridge: $812M --> $1.12B --> $1.5B",
                    fill_color=DARK_BG, font_size=20, font_color=WHITE, bold=True)

bridge = [
    ["Revenue Source", "FY26 Base", "FY27 Add", "FY27 Total", "FY28 Add", "FY28 Total"],
    ["Existing Account Growth", "$812M", "+$80M", "$892M", "+$100M", "$992M"],
    ["Migration Mega-Wave", "$0", "+$80M", "$80M", "+$120M", "$200M"],
    ["Agentic AI Factory", "$0", "+$50M", "$50M", "+$100M", "$150M"],
    ["New Diamond Activation", "$0", "+$60M", "$60M", "+$80M", "$140M"],
    ["Industry Data Clouds", "$0", "+$30M", "$30M", "+$50M", "$80M"],
    ["Ecosystem Flywheel (Po3)", "$0", "+$28M", "$28M", "+$42M", "$70M"],
    ["Risk / Attrition Buffer", "$0", "-$20M", "-$20M", "-$32M", "-$32M"],
    ["", "", "", "", "", ""],
    ["TOTAL", "$812M", "+$308M", "$1,120M", "+$380M", "$1,500M"],
]
add_table(slide, 200000, 700000, 8700000, 2700000, 10, 6, bridge)

add_text_box(slide, 200000, 3600000, 8700000, 300000,
             "SNOW Contribution Summary", 14, DARK_BG, True)

snow_summary = [
    ["Metric", "FY26", "FY27 Plan", "FY28 Plan"],
    ["TACV", "$72M", "$110M (+53%)", "$160M (+45%)"],
    ["Revenue Consumption", "~$24M", "$45M (+88%)", "$75M (+67%)"],
    ["Total SNOW Contribution", "$96M", "$155M", "$235M"],
    ["Use Cases", "218", "350 (+60%)", "500 (+43%)"],
]
add_table(slide, 200000, 3900000, 8700000, 1000000, 5, 4, snow_summary)

# ============================================================
# SLIDE 26: KEY PARTNERSHIP PRIORITIES FY27-28
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "FY27-28 Top Partnership Priorities",
                    fill_color=DARK_BG, font_size=20, font_color=WHITE, bold=True)

priorities = [
    ("1", "Sustain $1B+ and reach $1.5B by end of CY2028, becoming the first and largest partner alliance in the Data & AI industry.",
     SNOW_BLUE, 900000),
    ("2", "Scale Agentic AI: Deploy Cortex-powered agents across 150+ accounts via AI Refinery 2.0, targeting $250M pipeline.",
     ACCENT_GREEN, 1800000),
    ("3", "Industrialize Migrations: Execute $300M Migration Mega-Wave targeting Teradata, SAS, Oracle, Hadoop, and SAP BW workloads.",
     ACCENT_PURPLE, 2700000),
    ("4", "Deepen Ecosystem: Evolve Power of 3 into systematic flywheel with all ACN Business Groups (AWS, Azure, SAP, SFDC, GCP).",
     ACCENT_ORANGE, 3600000),
]

for num, text, color, y_pos in priorities:
    add_shape_with_text(slide, MSO_SHAPE.OVAL, 300000, y_pos, 500000, 500000,
                        num, fill_color=color, font_size=22, font_color=WHITE, bold=True)
    add_text_box(slide, 950000, y_pos + 50000, 7800000, 450000,
                 text, 12, RGBColor(0x33, 0x33, 0x33), False)

# ============================================================
# SLIDE 27: DIVIDER - GOVERNANCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 2000000, 9144000, 80000,
                    "", fill_color=RGBColor(0xEC, 0x48, 0x99))
add_text_box(slide, 500000, 2200000, 8000000, 700000,
             "Governance & Mobilization", 36, WHITE, True, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 2900000, 8000000, 400000,
             "Operating Rhythm, KPIs & Accountability Framework", 18, ACCENT_CYAN, False, PP_ALIGN.LEFT)
add_text_box(slide, 500000, 3600000, 4000000, 300000,
             "Section 10", 14, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 28: GOVERNANCE FRAMEWORK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Partnership Operating Rhythm & Governance",
                    fill_color=DARK_BG, font_size=20, font_color=WHITE, bold=True)

gov = [
    ["Cadence", "Forum", "Attendees", "Focus", "Deliverable"],
    ["Weekly", "Pipeline War Room", "Regional PDMs, Sales Leads", "Deal progression, blockers", "Pipeline scorecard"],
    ["Bi-Weekly", "Use Case Acceleration", "COE Pods, PSEs, Industry", "Use case deployment, go-lives", "Use case tracker"],
    ["Monthly", "Revenue Review", "GTM Leads, Finance, Ops", "Revenue tracking, forecast", "Monthly P&L"],
    ["Quarterly", "QBR + Exec Steering", "ELT Sponsors (both sides)", "Strategy, investments, escalations", "Exec dashboard"],
    ["Bi-Annual", "Strategic Planning Summit", "CEOs, CROs, CTOs", "Annual planning, vision", "Joint roadmap"],
    ["Annual", "Partnership Innovation Day", "Top 100 leaders + clients", "Innovation showcase, commitments", "Annual plan refresh"],
]
add_table(slide, 200000, 700000, 8700000, 1700000, 7, 5, gov)

add_text_box(slide, 200000, 2600000, 8700000, 300000,
             "KPI Scoreboard — Measured Weekly", 14, DARK_BG, True)

kpis = [
    ["KPI Category", "Metric", "FY27 Target", "Measurement", "Owner"],
    ["Revenue", "ACN Services Bookings", "$1,120M", "Monthly actuals vs. plan", "Global Alliance Lead"],
    ["Pipeline", "Qualified Pipeline Coverage", "3x target ($3.36B)", "Weekly pipeline review", "Regional GTM Leads"],
    ["TACV", "Snowflake TACV", "$110M", "Monthly TACV report", "PDM Team"],
    ["Activation", "Diamond Account Penetration", "57% (175/307)", "Quarterly activation report", "Field Sales"],
    ["Use Cases", "Go-Live Count", "350", "Monthly go-live tracker", "COE Leads"],
    ["Certifications", "Net New Certifications", "4,000", "Quarterly cert report", "Enablement"],
    ["NPS", "Joint Client Satisfaction", ">70 NPS", "Quarterly survey", "Marketing"],
]
add_table(slide, 200000, 2900000, 8700000, 2000000, 8, 5, kpis)

# ============================================================
# SLIDE 29: INVESTMENT ASK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 600000,
                    "Joint Investment Ask — FY27-28",
                    fill_color=DARK_BG, font_size=20, font_color=WHITE, bold=True)

investments = [
    ["Investment Area", "FY27 Ask", "FY28 Ask", "Total", "Expected ROI"],
    ["Agentic AI Factory (Solution Dev + GTM)", "$5M", "$8M", "$13M", "25x ($325M pipeline)"],
    ["Industry Data Cloud Development", "$4M", "$6M", "$10M", "20x ($200M pipeline)"],
    ["Migration Factory Tooling + Fund", "$3M", "$5M", "$8M", "37x ($300M pipeline)"],
    ["COE & Practice Scaling", "$6M", "$8M", "$14M", "N/A (capability build)"],
    ["Marketing & Demand Generation", "$3M", "$4M", "$7M", "15x ($105M pipeline)"],
    ["Executive Events & Innovation Labs", "$2M", "$3M", "$5M", "20x ($100M pipeline)"],
    ["Field Enablement & Certifications", "$2M", "$3M", "$5M", "N/A (capability build)"],
    ["TOTAL INVESTMENT", "$25M", "$37M", "$62M", "Overall: 22x ROI"],
]
add_table(slide, 200000, 700000, 8700000, 2200000, 9, 5, investments)

add_text_box(slide, 200000, 3100000, 8700000, 300000,
             "Funding Model", 14, DARK_BG, True)
add_text_box(slide, 200000, 3400000, 8700000, 1400000,
             "Proposed 50/50 split between Snowflake and Accenture, with milestone-based releases:\n"
             "  - 40% upfront on program approval\n"
             "  - 30% on H1 milestones (pipeline creation, use case go-lives, certifications)\n"
             "  - 30% on H2 milestones (revenue attainment, account activation)\n\n"
             "Additional leverage: Snowflake MDF, AWS/Azure co-funding for Power of 3 plays, "
             "ACN internal practice development funds, and client co-investment for custom solutions.",
             10, RGBColor(0x33, 0x33, 0x33), False)

# ============================================================
# SLIDE 30: CLOSING SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 0, 9144000, 80000,
                    "", fill_color=SNOW_BLUE)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 0, 5063500, 9144000, 80000,
                    "", fill_color=ACCENT_GREEN)

add_text_box(slide, 500000, 1000000, 8144000, 600000,
             "From $1B Milestone to Market Dominance", 32, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 500000, 1600000, 8144000, 500000,
             "The First Partnership to Define the\nData & AI Industry Together", 22, SNOW_BLUE, False, PP_ALIGN.CENTER)

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, 2500000, 2200000, 4144000, 30000,
                    "", fill_color=ACCENT_PURPLE)

items = [
    "$1.5B by CY2028",
    "200+ Diamond Accounts Activated",
    "500+ Use Cases Deployed",
    "12,000+ Certified Practitioners",
    "6 Industry Data Clouds",
    "150+ Agentic AI Deployments",
]
y = 2400000
for item in items:
    add_text_box(slide, 2000000, y, 5144000, 250000,
                 item, 14, ACCENT_CYAN, False, PP_ALIGN.CENTER)
    y += 280000

add_text_box(slide, 500000, 4500000, 8144000, 300000,
             "Confidential | March 2026 | Snowflake x Accenture", 12, MEDIUM_GRAY, False, PP_ALIGN.CENTER)


output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "FY27-28_ACN_SNOW_Path_to_1.5B_Business_Plan.pptx")
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
